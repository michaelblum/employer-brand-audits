from __future__ import annotations

import json
import sys
import tempfile
import unittest
from pathlib import Path

REPO_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(REPO_ROOT))

from scripts.publication_pipeline_fixture import (
    KILOS_METHODOLOGY_TEMPLATE_ID,
    generate_kilos_methodology_fixture,
)
from scripts.workbench_projection import project_audit_manifest


REFERENCE_KILOS_JSON = REPO_ROOT / "reference_publications" / "kilos-framework.json"
KILOS_MAPPING = REPO_ROOT / "reference_publications" / "Ontologies" / "KILOS" / "KILOS mapping.xlsx"
KILOS_INTRO = REPO_ROOT / "reference_publications" / "Ontologies" / "KILOS" / "KILOS Introduction.pptx"
KILOS_METHOD = REPO_ROOT / "reference_publications" / "Ontologies" / "KILOS" / "ST_ KILOS Methodology.pptx"
KILOS_TABLES = REPO_ROOT / "reference_publications" / "Ontologies" / "KILOS" / "KILOS tables.pptx"


class KilosMethodologyPipelineTests(unittest.TestCase):
    @unittest.skipUnless(REFERENCE_KILOS_JSON.exists(), "local KILOS references are not tracked")
    def test_local_kilos_references_are_readable_when_present(self) -> None:
        try:
            import openpyxl
            from pptx import Presentation
        except ImportError as exc:
            self.skipTest(f"reference readers unavailable: {exc}")

        reference = json.loads(REFERENCE_KILOS_JSON.read_text(encoding="utf-8"))
        tracked = json.loads((REPO_ROOT / "data" / "kilos-framework.json").read_text(encoding="utf-8"))
        workbook = openpyxl.load_workbook(KILOS_MAPPING, read_only=True, data_only=True)

        self.assertEqual(reference, tracked)
        self.assertEqual(workbook.sheetnames, ["Export"])
        self.assertEqual(workbook["Export"].calculate_dimension(force=True), "A1:D33")
        self.assertEqual(len(Presentation(KILOS_INTRO).slides), 2)
        self.assertEqual(len(Presentation(KILOS_METHOD).slides), 3)
        self.assertEqual(len(Presentation(KILOS_TABLES).slides), 13)

    def test_kilos_fixture_writes_manifest_and_required_views(self) -> None:
        with tempfile.TemporaryDirectory(dir=REPO_ROOT / "artifacts") as tmp:
            manifest_path = generate_kilos_methodology_fixture(Path(tmp) / "kilos")
            manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
            for relative_path in [
                "pipeline-intake.json",
                "ontology-source-roster.json",
                "kilos-browser.json",
                "mapping-workbook.json",
                "methodology-deck.json",
                "scorecard-tables.json",
                "report-section-snippets.html",
                "l4-publication.html",
            ]:
                self.assertTrue((manifest_path.parent / relative_path).exists(), relative_path)

        self.assertEqual(manifest["template_id"], KILOS_METHODOLOGY_TEMPLATE_ID)
        self.assertNotIn("reference_publications", json.dumps(manifest))

    def test_kilos_fixture_preserves_ontology_and_mapping_counts(self) -> None:
        with tempfile.TemporaryDirectory(dir=REPO_ROOT / "artifacts") as tmp:
            manifest_path = generate_kilos_methodology_fixture(Path(tmp) / "kilos")
            root = manifest_path.parent
            sources = json.loads((root / "ontology-source-roster.json").read_text(encoding="utf-8"))
            browser = json.loads((root / "kilos-browser.json").read_text(encoding="utf-8"))
            mapping = json.loads((root / "mapping-workbook.json").read_text(encoding="utf-8"))
            deck = json.loads((root / "methodology-deck.json").read_text(encoding="utf-8"))
            scorecards = json.loads((root / "scorecard-tables.json").read_text(encoding="utf-8"))

        self.assertEqual(len(sources["ontology_sources"]), 5)
        self.assertEqual(browser["framework_id"], "KILOS")
        self.assertEqual(browser["framework_version"], "1.0")
        self.assertEqual(len(browser["pillars"]), 5)
        self.assertEqual(len(browser["factors"]), 29)
        self.assertEqual(browser["survey_label_assignments"], 33)
        self.assertEqual(mapping["sheet_name"], "Export")
        self.assertEqual(mapping["row_count"], 33)
        self.assertEqual(mapping["response_rows"], 32)
        self.assertEqual(mapping["kilos_rows"], 31)
        self.assertEqual(mapping["non_kilos_rows"], 1)
        self.assertEqual(deck["total_slides"], 18)
        self.assertEqual(scorecards["tables_deck_slides"], 13)
        self.assertEqual(scorecards["tables"], 11)
        self.assertEqual(scorecards["pictures"], 25)

    def test_kilos_projection_groups_l4_view(self) -> None:
        with tempfile.TemporaryDirectory(dir=REPO_ROOT / "artifacts") as tmp:
            manifest_path = generate_kilos_methodology_fixture(Path(tmp) / "kilos")
            projection = project_audit_manifest(manifest_path)

        self.assertIn("publication.l4_publication.bundle", {group["slot"] for group in projection["artifact_groups"]})

    def test_kilos_fixture_is_registered_in_command_surface(self) -> None:
        from scripts.eba_cli import FIXTURE_GENERATORS, demo_recipe_lines, validation_commands

        self.assertIn("kilos-methodology", FIXTURE_GENERATORS)
        self.assertTrue(
            any(command[-1] == "tests/test_publication_kilos_methodology.py" for command in validation_commands())
        )
        with tempfile.TemporaryDirectory(dir=REPO_ROOT / "artifacts") as tmp:
            manifest = Path(tmp) / "manifest.json"
            manifest.write_text(json.dumps({"template_id": KILOS_METHODOLOGY_TEMPLATE_ID}), encoding="utf-8")
            lines = demo_recipe_lines(fixture=None, manifest=manifest)

        self.assertIn("KILOS Browser", "\n".join(lines))
        self.assertIn("Mapping Workbook", "\n".join(lines))


if __name__ == "__main__":
    unittest.main()
